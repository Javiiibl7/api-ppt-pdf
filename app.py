from fastapi import FastAPI, HTTPException
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import uvicorn

app = FastAPI()

# 📌 Función para crear PowerPoint
def create_powerpoint(images, output_pptx):
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Diapositiva en blanco

        for i in range(0, len(images), 4):
            slide = prs.slides.add_slide(slide_layout)
            img_positions = [(Inches(0.5), Inches(0.5)), (Inches(5.5), Inches(0.5)),
                             (Inches(0.5), Inches(4.5)), (Inches(5.5), Inches(4.5))]

            for j, img_path in enumerate(images[i:i+4]):
                left, top = img_positions[j]
                img = Image.open(img_path)
                width, height = img.size
                aspect_ratio = width / height
                img_width = Inches(4) if aspect_ratio > 1 else Inches(3)
                img_height = Inches(3) if aspect_ratio < 1 else Inches(4)
                slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

        prs.save(output_pptx)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al crear PowerPoint: {str(e)}")

# 📌 Función para crear PDF
def create_pdf(images, output_pdf):
    try:
        c = canvas.Canvas(output_pdf, pagesize=letter)
        width, height = letter

        for i in range(0, len(images), 4):
            for j in range(4):
                if i + j < len(images):
                    img_path = images[i + j]
                    top = height - (j // 2) * 280 - 40
                    left = (j % 2) * 300 + 40
                    c.drawImage(img_path, left, top, width=250, height=180)
            c.showPage()

        c.save()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al crear PDF: {str(e)}")

# 📌 Endpoint de la API para generar PowerPoint y PDF
@app.post("/generate_ppt_pdf")
def generate_ppt_pdf(folder_path: str):
    """
    Recibe la ruta de una carpeta y genera un PowerPoint y un PDF con imágenes dentro de la carpeta.
    """
    if not os.path.exists(folder_path):
        raise HTTPException(status_code=400, detail="La carpeta especificada no existe.")

    images = [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

    if not images:
        raise HTTPException(status_code=404, detail="No se encontraron imágenes en la carpeta.")

    output_pptx = os.path.join(folder_path, "output_presentation.pptx")
    output_pdf = os.path.join(folder_path, "output_document.pdf")

    create_powerpoint(images, output_pptx)
    create_pdf(images, output_pdf)

    return {"pptx_file": output_pptx, "pdf_file": output_pdf}

# 📌 Endpoint raíz para evitar error 404
@app.get("/")
def home():
    return {"message": "API de generación de PowerPoint y PDF funcionando correctamente"}

# 📌 Ajuste del puerto dinámico para Render
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))  # Usa el puerto que Render asigna
    uvicorn.run(app, host="0.0.0.0", port=port)
